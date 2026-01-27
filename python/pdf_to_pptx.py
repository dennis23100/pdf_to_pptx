#!/usr/bin/env python3
"""
📄 PDF/圖片 轉 PPTX 可編輯工具
- OCR 偵測文字位置
- 智能背景色覆蓋原文字
- 生成可編輯的 PPTX 文字方塊

作者: Claude AI
授權: MIT License
"""

import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import pytesseract
from pathlib import Path
import argparse
import sys
import os

# ============================================================
# 輔助函數
# ============================================================

def clamp_color(c):
    """確保顏色值在 0-255 範圍內"""
    return max(0, min(255, int(c)))


def get_background_color(image, bbox, margin=5):
    """獲取文字周圍的背景顏色"""
    x, y, w, h = bbox
    img_h, img_w = image.shape[:2]
    
    regions = []
    
    # 取四周區域
    if y > margin:
        regions.append(image[max(0, y-margin):y, x:x+w])
    if y + h + margin < img_h:
        regions.append(image[y+h:min(img_h, y+h+margin), x:x+w])
    if x > margin:
        regions.append(image[y:y+h, max(0, x-margin):x])
    if x + w + margin < img_w:
        regions.append(image[y:y+h, x+w:min(img_w, x+w+margin)])
    
    if regions:
        valid = [r for r in regions if r.size > 0]
        if valid:
            all_pixels = np.vstack([r.reshape(-1, 3) for r in valid])
            if len(all_pixels) > 0:
                median = np.median(all_pixels, axis=0)
                # BGR to RGB
                return (clamp_color(median[2]), clamp_color(median[1]), clamp_color(median[0]))
    
    return (245, 240, 230)  # 預設米色


def get_text_color(image, bbox):
    """獲取文字顏色（取最暗的顏色）"""
    x, y, w, h = bbox
    
    # 確保在範圍內
    x = max(0, x)
    y = max(0, y)
    x2 = min(image.shape[1], x + w)
    y2 = min(image.shape[0], y + h)
    
    roi = image[y:y2, x:x2]
    
    if roi.size == 0:
        return (50, 50, 50)
    
    pixels = roi.reshape(-1, 3)
    brightness = np.sum(pixels, axis=1)
    threshold = np.percentile(brightness, 20)
    dark_pixels = pixels[brightness <= threshold]
    
    if len(dark_pixels) > 0:
        median = np.median(dark_pixels, axis=0)
        # BGR to RGB
        return (clamp_color(median[2]), clamp_color(median[1]), clamp_color(median[0]))
    
    return (50, 50, 50)


# ============================================================
# 核心功能
# ============================================================

def detect_text(image, lang='chi_tra+eng', conf_threshold=25):
    """
    使用 Tesseract OCR 偵測文字
    
    Args:
        image: OpenCV 圖片 (BGR)
        lang: OCR 語言
        conf_threshold: 信心度閾值
    
    Returns:
        list: 文字區域列表
    """
    # 轉換為 PIL Image
    pil_image = Image.fromarray(cv2.cvtColor(image, cv2.COLOR_BGR2RGB))
    
    # OCR
    data = pytesseract.image_to_data(pil_image, lang=lang, output_type=pytesseract.Output.DICT)
    
    words = []
    for i in range(len(data['text'])):
        text = data['text'][i].strip()
        conf = int(data['conf'][i])
        
        if conf > conf_threshold and text:
            words.append({
                'text': text,
                'bbox': (data['left'][i], data['top'][i], data['width'][i], data['height'][i]),
                'conf': conf,
                'line': data['line_num'][i],
                'block': data['block_num'][i]
            })
    
    return words


def group_words_to_lines(words):
    """將單詞按行分組"""
    lines = {}
    
    for w in words:
        key = (w['block'], w['line'])
        if key not in lines:
            lines[key] = []
        lines[key].append(w)
    
    # 合併每一行
    result = []
    for key, line_words in lines.items():
        line_words.sort(key=lambda x: x['bbox'][0])
        
        text = ' '.join([w['text'] for w in line_words])
        x_min = min(w['bbox'][0] for w in line_words)
        y_min = min(w['bbox'][1] for w in line_words)
        x_max = max(w['bbox'][0] + w['bbox'][2] for w in line_words)
        y_max = max(w['bbox'][1] + w['bbox'][3] for w in line_words)
        
        avg_h = sum(w['bbox'][3] for w in line_words) / len(line_words)
        
        result.append({
            'text': text,
            'bbox': (x_min, y_min, x_max - x_min, y_max - y_min),
            'font_size': int(avg_h * 0.72)
        })
    
    return result


def process_image(image, lang='chi_tra+eng', conf_threshold=25):
    """
    處理單張圖片
    
    Returns:
        tuple: (原圖, 合併後的行列表)
    """
    print(f"  🔍 OCR 偵測中...")
    words = detect_text(image, lang, conf_threshold)
    print(f"     找到 {len(words)} 個文字")
    
    lines = group_words_to_lines(words)
    print(f"     合併為 {len(lines)} 行")
    
    # 為每行計算背景色和文字色
    for line in lines:
        line['bg_color'] = get_background_color(image, line['bbox'])
        line['text_color'] = get_text_color(image, line['bbox'])
    
    return image, lines


def create_slide(prs, image, lines):
    """創建投影片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版面
    
    img_h, img_w = image.shape[:2]
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    
    # 保存原圖作為背景
    temp_path = '/tmp/pptx_bg_temp.png'
    cv2.imwrite(temp_path, image)
    
    # 添加背景圖片
    slide.shapes.add_picture(
        temp_path, Inches(0), Inches(0),
        width=slide_w, height=slide_h
    )
    
    # 座標轉換比例
    scale_x = slide_w.inches / img_w
    scale_y = slide_h.inches / img_h
    
    # 為每行添加覆蓋層和文字方塊
    for line in lines:
        x, y, w, h = line['bbox']
        
        left = Inches(x * scale_x)
        top = Inches(y * scale_y)
        width = Inches(w * scale_x)
        height = Inches(h * scale_y)
        
        # 1. 添加背景色矩形覆蓋原文字
        cover = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        
        bg = line['bg_color']
        cover.fill.solid()
        cover.fill.fore_color.rgb = RGBColor(bg[0], bg[1], bg[2])
        cover.line.fill.background()  # 無邊框
        
        # 2. 添加可編輯文字方塊
        textbox = slide.shapes.add_textbox(
            left, top, width, Inches(h * scale_y * 1.3)
        )
        tf = textbox.text_frame
        tf.word_wrap = False
        
        p = tf.paragraphs[0]
        p.text = line['text']
        
        font_size = max(8, min(line['font_size'], 48))
        p.font.size = Pt(font_size)
        p.font.name = 'Microsoft JhengHei'  # 微軟正黑體
        
        tc = line['text_color']
        p.font.color.rgb = RGBColor(tc[0], tc[1], tc[2])
        
        textbox.fill.background()  # 透明背景
    
    # 清理臨時檔案
    if os.path.exists(temp_path):
        os.remove(temp_path)


def convert_to_pptx(input_path, output_path, lang='chi_tra+eng', dpi=200, conf=25):
    """
    主轉換函數
    
    Args:
        input_path: 輸入檔案路徑 (PDF 或圖片)
        output_path: 輸出 PPTX 路徑
        lang: OCR 語言
        dpi: PDF 轉換 DPI
        conf: OCR 信心度閾值
    """
    input_path = Path(input_path)
    
    # 讀取圖片
    if input_path.suffix.lower() == '.pdf':
        try:
            from pdf2image import convert_from_path
            print(f"📄 轉換 PDF (DPI: {dpi})...")
            pil_images = convert_from_path(str(input_path), dpi=dpi)
            images = [cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR) for img in pil_images]
        except ImportError:
            print("❌ 錯誤: 請安裝 pdf2image: pip install pdf2image")
            print("   還需要安裝 poppler: ")
            print("   - Windows: 下載 https://github.com/oschwartz10612/poppler-windows/releases")
            print("   - Mac: brew install poppler")
            print("   - Linux: sudo apt install poppler-utils")
            sys.exit(1)
    else:
        image = cv2.imread(str(input_path))
        if image is None:
            print(f"❌ 錯誤: 無法讀取圖片 {input_path}")
            sys.exit(1)
        images = [image]
    
    print(f"📊 共 {len(images)} 頁")
    
    # 創建 PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for i, image in enumerate(images):
        print(f"\n📖 處理第 {i+1}/{len(images)} 頁...")
        
        image, lines = process_image(image, lang, conf)
        
        print(f"  📊 創建投影片...")
        create_slide(prs, image, lines)
    
    # 保存
    prs.save(str(output_path))
    print(f"\n✅ 完成！輸出: {output_path}")


# ============================================================
# 命令列介面
# ============================================================

def main():
    parser = argparse.ArgumentParser(
        description='📄 PDF/圖片 轉 PPTX 可編輯工具',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
範例:
  python pdf_to_pptx.py input.pdf output.pptx
  python pdf_to_pptx.py image.png output.pptx --lang eng
  python pdf_to_pptx.py input.pdf output.pptx --lang chi_tra+eng --dpi 300

支援的語言代碼:
  chi_tra     繁體中文
  chi_sim     簡體中文
  eng         英文
  jpn         日文
  kor         韓文
  
  多語言: chi_tra+eng (繁中+英文)
        """
    )
    
    parser.add_argument('input', help='輸入檔案 (PDF 或圖片)')
    parser.add_argument('output', help='輸出 PPTX 檔案')
    parser.add_argument('--lang', default='chi_tra+eng', help='OCR 語言 (預設: chi_tra+eng)')
    parser.add_argument('--dpi', type=int, default=200, help='PDF 轉換 DPI (預設: 200)')
    parser.add_argument('--conf', type=int, default=25, help='OCR 信心度閾值 (預設: 25)')
    parser.add_argument('--preview', action='store_true', help='生成預覽圖片')
    
    args = parser.parse_args()
    
    # 檢查輸入檔案
    if not Path(args.input).exists():
        print(f"❌ 錯誤: 找不到輸入檔案 {args.input}")
        sys.exit(1)
    
    # 轉換
    convert_to_pptx(args.input, args.output, args.lang, args.dpi, args.conf)
    
    # 生成預覽
    if args.preview:
        try:
            preview_path = Path(args.output).with_suffix('.preview.png')
            image = cv2.imread(args.input) if not args.input.lower().endswith('.pdf') else None
            if image is not None:
                cv2.imwrite(str(preview_path), image)
                print(f"📸 預覽: {preview_path}")
        except Exception as e:
            print(f"⚠️ 無法生成預覽: {e}")


if __name__ == '__main__':
    main()
